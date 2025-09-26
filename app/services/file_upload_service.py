"""
文件上传与处理服务 - 支持多种文件格式
用于研究方向确定和文献导入
"""

import os
import json
import asyncio
from typing import List, Dict, Optional, Any, BinaryIO
from pathlib import Path
import tempfile
import shutil
from datetime import datetime
from loguru import logger
import aiofiles
import magic
import hashlib
from sqlalchemy.orm import Session

# 文档处理库
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup

from app.core.config import settings
from app.services.ai_service import AIService
from app.services.pdf_processor import PDFProcessor
from app.models.user import User
from app.models.project import Project


class FileUploadService:
    """文件上传与处理服务"""
    
    def __init__(self, db: Session):
        self.db = db
        self.ai_service = AIService()
        self.pdf_processor = PDFProcessor()
        
        # 支持的文件类型
        self.supported_types = {
            # 文档类型
            "application/pdf": {"ext": ".pdf", "category": "document"},
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": {"ext": ".docx", "category": "document"},
            "application/msword": {"ext": ".doc", "category": "document"},
            "text/plain": {"ext": ".txt", "category": "document"},
            
            # 演示文稿
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": {"ext": ".pptx", "category": "presentation"},
            "application/vnd.ms-powerpoint": {"ext": ".ppt", "category": "presentation"},
            
            # 表格文件
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": {"ext": ".xlsx", "category": "spreadsheet"},
            "application/vnd.ms-excel": {"ext": ".xls", "category": "spreadsheet"},
            "text/csv": {"ext": ".csv", "category": "spreadsheet"},
            
            # Zotero导出
            "application/json": {"ext": ".json", "category": "bibliography"},
            "text/xml": {"ext": ".xml", "category": "bibliography"},
            "application/xml": {"ext": ".xml", "category": "bibliography"},
            "text/html": {"ext": ".html", "category": "bibliography"},
            
            # 其他格式
            "application/rtf": {"ext": ".rtf", "category": "document"},
            "text/markdown": {"ext": ".md", "category": "document"}
        }
        
        # 文件大小限制（MB）
        self.size_limits = {
            "document": 50,  # 50MB
            "presentation": 100,  # 100MB
            "spreadsheet": 20,  # 20MB
            "bibliography": 10  # 10MB
        }
    
    async def upload_and_analyze_file(
        self,
        file_content: bytes,
        filename: str,
        user: User,
        analysis_type: str = "research_direction",
        progress_callback = None
    ) -> Dict:
        """
        上传并分析文件
        
        Args:
            file_content: 文件内容
            filename: 文件名
            user: 用户对象
            analysis_type: 分析类型 (research_direction, literature_import)
            progress_callback: 进度回调函数
            
        Returns:
            分析结果
        """
        try:
            logger.info(f"开始处理文件: {filename}, 用户: {user.username}")
            
            if progress_callback:
                await progress_callback("验证文件格式", 5, {"filename": filename})
            
            # 第一步：文件验证
            validation_result = await self._validate_file(file_content, filename)
            if not validation_result["valid"]:
                return {"success": False, "error": validation_result["error"]}
            
            file_info = validation_result["file_info"]
            
            if progress_callback:
                await progress_callback("保存文件", 15, file_info)
            
            # 第二步：保存文件
            file_path = await self._save_file(file_content, filename, user.id)
            
            if progress_callback:
                await progress_callback("提取文件内容", 30, {"file_path": file_path})
            
            # 第三步：提取文件内容
            content_result = await self._extract_file_content(file_path, file_info)
            if not content_result["success"]:
                return {"success": False, "error": content_result["error"]}
            
            extracted_content = content_result["content"]
            
            if progress_callback:
                await progress_callback("AI分析内容", 60, {"content_length": len(extracted_content)})
            
            # 第四步：AI分析
            if analysis_type == "research_direction":
                analysis_result = await self._analyze_research_direction(extracted_content, filename)
            elif analysis_type == "literature_import":
                analysis_result = await self._analyze_literature_import(extracted_content, filename)
            else:
                return {"success": False, "error": f"不支持的分析类型: {analysis_type}"}
            
            if progress_callback:
                await progress_callback("分析完成", 100, {"analysis_type": analysis_type})
            
            # 第五步：清理临时文件
            try:
                os.unlink(file_path)
            except:
                pass
            
            return {
                "success": True,
                "file_info": file_info,
                "extracted_content": extracted_content[:1000],  # 返回部分内容用于展示
                "analysis_result": analysis_result
            }
            
        except Exception as e:
            logger.error(f"文件处理失败: {e}")
            return {"success": False, "error": str(e)}
    
    async def batch_upload_literature(
        self,
        files_data: List[Dict],  # [{"content": bytes, "filename": str}, ...]
        user: User,
        project: Optional[Project] = None,
        progress_callback = None
    ) -> Dict:
        """
        批量上传文献文件
        
        Args:
            files_data: 文件数据列表
            user: 用户对象
            project: 项目对象（可选）
            progress_callback: 进度回调函数
            
        Returns:
            批量处理结果
        """
        try:
            total_files = len(files_data)
            logger.info(f"开始批量处理 {total_files} 个文件")
            
            if progress_callback:
                await progress_callback("开始批量处理", 0, {"total_files": total_files})
            
            results = {
                "success": True,
                "total_files": total_files,
                "processed_files": 0,
                "successful_files": 0,
                "failed_files": 0,
                "file_results": [],
                "extracted_literature": []
            }
            
            # 并行处理文件（限制并发数）
            semaphore = asyncio.Semaphore(3)  # 最多3个文件同时处理
            
            async def process_single_file(file_data: Dict, index: int):
                async with semaphore:
                    try:
                        result = await self.upload_and_analyze_file(
                            file_data["content"],
                            file_data["filename"],
                            user,
                            "literature_import"
                        )
                        
                        results["processed_files"] += 1
                        
                        if result["success"]:
                            results["successful_files"] += 1
                            if "literature_data" in result.get("analysis_result", {}):
                                results["extracted_literature"].extend(
                                    result["analysis_result"]["literature_data"]
                                )
                        else:
                            results["failed_files"] += 1
                        
                        results["file_results"].append({
                            "filename": file_data["filename"],
                            "success": result["success"],
                            "error": result.get("error"),
                            "literature_count": len(result.get("analysis_result", {}).get("literature_data", []))
                        })
                        
                        # 更新进度
                        if progress_callback:
                            progress = (results["processed_files"] / total_files) * 100
                            await progress_callback(
                                f"处理文件 {results['processed_files']}/{total_files}",
                                progress,
                                {
                                    "processed": results["processed_files"],
                                    "successful": results["successful_files"],
                                    "failed": results["failed_files"]
                                }
                            )
                        
                        return result
                        
                    except Exception as e:
                        logger.error(f"处理文件 {file_data['filename']} 失败: {e}")
                        results["processed_files"] += 1
                        results["failed_files"] += 1
                        results["file_results"].append({
                            "filename": file_data["filename"],
                            "success": False,
                            "error": str(e),
                            "literature_count": 0
                        })
                        return {"success": False, "error": str(e)}
            
            # 创建并执行任务
            tasks = [
                process_single_file(file_data, i) 
                for i, file_data in enumerate(files_data)
            ]
            
            await asyncio.gather(*tasks, return_exceptions=True)
            
            if progress_callback:
                await progress_callback("批量处理完成", 100, results)
            
            logger.info(f"批量处理完成 - 成功: {results['successful_files']}, 失败: {results['failed_files']}")
            
            return results
            
        except Exception as e:
            logger.error(f"批量文件处理失败: {e}")
            return {"success": False, "error": str(e)}
    
    async def _validate_file(self, file_content: bytes, filename: str) -> Dict:
        """验证文件"""
        try:
            # 检查文件大小
            file_size_mb = len(file_content) / (1024 * 1024)
            
            # 检测MIME类型
            mime_type = magic.from_buffer(file_content, mime=True)
            
            if mime_type not in self.supported_types:
                return {
                    "valid": False,
                    "error": f"不支持的文件类型: {mime_type}"
                }
            
            file_info = self.supported_types[mime_type]
            category = file_info["category"]
            
            # 检查文件大小限制
            if file_size_mb > self.size_limits[category]:
                return {
                    "valid": False,
                    "error": f"文件过大: {file_size_mb:.1f}MB，限制: {self.size_limits[category]}MB"
                }
            
            # 计算文件哈希
            file_hash = hashlib.md5(file_content).hexdigest()
            
            return {
                "valid": True,
                "file_info": {
                    "filename": filename,
                    "mime_type": mime_type,
                    "category": category,
                    "size_mb": file_size_mb,
                    "hash": file_hash,
                    "extension": file_info["ext"]
                }
            }
            
        except Exception as e:
            logger.error(f"文件验证失败: {e}")
            return {"valid": False, "error": f"文件验证失败: {str(e)}"}
    
    async def _save_file(self, file_content: bytes, filename: str, user_id: int) -> str:
        """保存文件到临时目录"""
        try:
            # 创建用户临时目录
            temp_dir = Path(tempfile.gettempdir()) / "research_platform" / str(user_id)
            temp_dir.mkdir(parents=True, exist_ok=True)
            
            # 生成唯一文件名
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_filename = "".join(c for c in filename if c.isalnum() or c in "._-")
            temp_filename = f"{timestamp}_{safe_filename}"
            file_path = temp_dir / temp_filename
            
            # 保存文件
            async with aiofiles.open(file_path, "wb") as f:
                await f.write(file_content)
            
            return str(file_path)
            
        except Exception as e:
            logger.error(f"保存文件失败: {e}")
            raise
    
    async def _extract_file_content(self, file_path: str, file_info: Dict) -> Dict:
        """提取文件内容"""
        try:
            category = file_info["category"]
            mime_type = file_info["mime_type"]
            
            if category == "document":
                if mime_type == "application/pdf":
                    content = await self._extract_pdf_content(file_path)
                elif "word" in mime_type or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    content = await self._extract_docx_content(file_path)
                elif mime_type == "text/plain":
                    content = await self._extract_text_content(file_path)
                else:
                    content = ""
                    
            elif category == "presentation":
                content = await self._extract_pptx_content(file_path)
                
            elif category == "spreadsheet":
                content = await self._extract_excel_content(file_path)
                
            elif category == "bibliography":
                content = await self._extract_bibliography_content(file_path, mime_type)
                
            else:
                content = ""
            
            if not content:
                return {"success": False, "error": "无法提取文件内容"}
            
            return {"success": True, "content": content}
            
        except Exception as e:
            logger.error(f"提取文件内容失败: {e}")
            return {"success": False, "error": str(e)}
    
    async def _extract_pdf_content(self, file_path: str) -> str:
        """提取PDF内容"""
        try:
            # 首先尝试使用MinerU
            content = await self.pdf_processor.extract_text(file_path)
            if content and len(content.strip()) > 100:
                return content
            
            # 如果MinerU失败，使用PyPDF2作为备用
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content = ""
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
                return content.strip()
                
        except Exception as e:
            logger.error(f"PDF内容提取失败: {e}")
            return ""
    
    async def _extract_docx_content(self, file_path: str) -> str:
        """提取DOCX内容"""
        try:
            doc = docx.Document(file_path)
            content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))
            
            return "\n".join(content)
            
        except Exception as e:
            logger.error(f"DOCX内容提取失败: {e}")
            return ""
    
    async def _extract_pptx_content(self, file_path: str) -> str:
        """提取PPTX内容"""
        try:
            prs = Presentation(file_path)
            content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== 幻灯片 {slide_num} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                content.extend(slide_content)
            
            return "\n".join(content)
            
        except Exception as e:
            logger.error(f"PPTX内容提取失败: {e}")
            return ""
    
    async def _extract_excel_content(self, file_path: str) -> str:
        """提取Excel内容"""
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path, sheet_name=None)
                if isinstance(df, dict):
                    # 多个工作表，合并内容
                    all_content = []
                    for sheet_name, sheet_df in df.items():
                        all_content.append(f"=== 工作表: {sheet_name} ===")
                        all_content.append(sheet_df.to_string(index=False))
                    return "\n".join(all_content)
                else:
                    df = df
            
            return df.to_string(index=False)
            
        except Exception as e:
            logger.error(f"Excel内容提取失败: {e}")
            return ""
    
    async def _extract_text_content(self, file_path: str) -> str:
        """提取纯文本内容"""
        try:
            async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                return await f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                async with aiofiles.open(file_path, "r", encoding="gbk") as f:
                    return await f.read()
            except:
                async with aiofiles.open(file_path, "r", encoding="latin-1") as f:
                    return await f.read()
        except Exception as e:
            logger.error(f"文本内容提取失败: {e}")
            return ""
    
    async def _extract_bibliography_content(self, file_path: str, mime_type: str) -> str:
        """提取参考文献内容"""
        try:
            if mime_type == "application/json":
                async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                    content = await f.read()
                    return content
                    
            elif "xml" in mime_type:
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
                    
            elif "html" in mime_type:
                with open(file_path, "r", encoding="utf-8") as f:
                    html_content = f.read()
                    soup = BeautifulSoup(html_content, "html.parser")
                    return soup.get_text()
            
            return ""
            
        except Exception as e:
            logger.error(f"参考文献内容提取失败: {e}")
            return ""
    
    async def _analyze_research_direction(self, content: str, filename: str) -> Dict:
        """分析研究方向"""
        try:
            analysis_prompt = f"""
请分析以下文档内容，提取研究方向和关键信息：

文档名称: {filename}
文档内容:
{content[:4000]}

请提取以下信息，以JSON格式返回：
{{
    "research_direction": "主要研究方向",
    "keywords": ["关键词1", "关键词2", "关键词3", ...],
    "research_field": "所属学科领域",
    "research_objectives": ["研究目标1", "研究目标2", ...],
    "methodology": ["研究方法1", "研究方法2", ...],
    "key_concepts": ["核心概念1", "核心概念2", ...],
    "application_areas": ["应用领域1", "应用领域2", ...],
    "confidence_score": 8.5
}}

要求：
1. 准确提取核心研究内容
2. 关键词要具体且相关性强
3. 研究方向要明确具体
4. 置信度评分（1-10分）
"""
            
            response = await self.ai_service.generate_completion(
                analysis_prompt,
                model="gpt-4",
                max_tokens=800,
                temperature=0.2
            )
            
            if response.get("success"):
                try:
                    analysis_result = json.loads(response["content"])
                    return analysis_result
                except json.JSONDecodeError:
                    # 如果JSON解析失败，返回基本信息
                    return {
                        "research_direction": "文档分析",
                        "keywords": ["研究", "分析"],
                        "research_field": "综合",
                        "confidence_score": 5.0
                    }
            
            return {"error": "AI分析失败"}
            
        except Exception as e:
            logger.error(f"研究方向分析失败: {e}")
            return {"error": str(e)}
    
    async def _analyze_literature_import(self, content: str, filename: str) -> Dict:
        """分析文献导入"""
        try:
            # 检测是否为Zotero导出格式
            if self._is_zotero_format(content):
                return await self._parse_zotero_export(content)
            
            # 检测是否为文献列表
            if self._is_literature_list(content):
                return await self._parse_literature_list(content)
            
            # 通用文献内容分析
            return await self._analyze_general_literature(content, filename)
            
        except Exception as e:
            logger.error(f"文献导入分析失败: {e}")
            return {"error": str(e)}
    
    def _is_zotero_format(self, content: str) -> bool:
        """检测是否为Zotero导出格式"""
        try:
            if content.strip().startswith('{') or content.strip().startswith('['):
                data = json.loads(content)
                if isinstance(data, list) and len(data) > 0:
                    first_item = data[0]
                    # 检查Zotero典型字段
                    zotero_fields = ["title", "creators", "abstractNote", "DOI", "url"]
                    return any(field in first_item for field in zotero_fields)
            return False
        except:
            return False
    
    def _is_literature_list(self, content: str) -> bool:
        """检测是否为文献列表"""
        lines = content.strip().split('\n')
        # 简单启发式：如果有多行且包含年份模式
        if len(lines) > 3:
            year_pattern_count = sum(1 for line in lines if any(str(year) in line for year in range(1990, 2025)))
            return year_pattern_count > len(lines) * 0.3
        return False
    
    async def _parse_zotero_export(self, content: str) -> Dict:
        """解析Zotero导出内容"""
        try:
            data = json.loads(content)
            literature_data = []
            
            for item in data:
                lit_item = {
                    "title": item.get("title", ""),
                    "authors": [],
                    "abstract": item.get("abstractNote", ""),
                    "year": item.get("date", "")[:4] if item.get("date") else None,
                    "journal": item.get("publicationTitle", ""),
                    "doi": item.get("DOI", ""),
                    "url": item.get("url", ""),
                    "source": "zotero_import"
                }
                
                # 处理作者信息
                if "creators" in item:
                    for creator in item["creators"]:
                        name = f"{creator.get('firstName', '')} {creator.get('lastName', '')}".strip()
                        if name:
                            lit_item["authors"].append({"name": name})
                
                if lit_item["title"]:
                    literature_data.append(lit_item)
            
            return {
                "literature_data": literature_data,
                "format": "zotero",
                "count": len(literature_data)
            }
            
        except Exception as e:
            logger.error(f"Zotero导出解析失败: {e}")
            return {"error": str(e)}
    
    async def _parse_literature_list(self, content: str) -> Dict:
        """解析文献列表"""
        try:
            # 使用AI解析文献列表
            parse_prompt = f"""
请解析以下文献列表，提取每篇文献的信息：

文献列表:
{content[:3000]}

请以JSON格式返回文献信息：
{{
    "literature_data": [
        {{
            "title": "文献标题",
            "authors": [{{"name": "作者名"}}, ...],
            "journal": "期刊名",
            "year": 2023,
            "doi": "DOI号",
            "abstract": "摘要（如果有）"
        }},
        ...
    ]
}}

要求：
1. 尽可能准确提取信息
2. 如果信息不完整，设为空字符串或null
3. 年份转换为数字
"""
            
            response = await self.ai_service.generate_completion(
                parse_prompt,
                model="gpt-4",
                max_tokens=2000,
                temperature=0.1
            )
            
            if response.get("success"):
                try:
                    result = json.loads(response["content"])
                    literature_data = result.get("literature_data", [])
                    
                    # 为每篇文献添加来源标记
                    for lit in literature_data:
                        lit["source"] = "list_import"
                    
                    return {
                        "literature_data": literature_data,
                        "format": "list",
                        "count": len(literature_data)
                    }
                except json.JSONDecodeError:
                    pass
            
            return {"error": "文献列表解析失败"}
            
        except Exception as e:
            logger.error(f"文献列表解析失败: {e}")
            return {"error": str(e)}
    
    async def _analyze_general_literature(self, content: str, filename: str) -> Dict:
        """分析通用文献内容"""
        try:
            analysis_prompt = f"""
请分析以下文献内容，提取关键信息：

文件名: {filename}
内容:
{content[:2000]}

请提取信息并以JSON格式返回：
{{
    "title": "推测的文献标题",
    "research_topics": ["研究主题1", "研究主题2", ...],
    "keywords": ["关键词1", "关键词2", ...],
    "methodology": ["研究方法1", "研究方法2", ...],
    "key_findings": ["主要发现1", "主要发现2", ...],
    "literature_type": "期刊论文/会议论文/书籍/报告",
    "confidence_score": 7.5
}}

要求：
1. 基于内容推测文献信息
2. 提取最相关的关键词和主题
3. 评估分析置信度（1-10分）
"""
            
            response = await self.ai_service.generate_completion(
                analysis_prompt,
                model="gpt-4",
                max_tokens=600,
                temperature=0.3
            )
            
            if response.get("success"):
                try:
                    analysis_result = json.loads(response["content"])
                    return {
                        "content_analysis": analysis_result,
                        "format": "general",
                        "extractable_info": True
                    }
                except json.JSONDecodeError:
                    pass
            
            return {
                "content_analysis": {"title": filename, "keywords": []},
                "format": "general",
                "extractable_info": False
            }
            
        except Exception as e:
            logger.error(f"通用文献分析失败: {e}")
            return {"error": str(e)}